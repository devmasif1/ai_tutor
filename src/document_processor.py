import os
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_community.document_loaders import PyPDFLoader
from langchain_experimental.text_splitter import SemanticChunker
from langchain_openai import OpenAIEmbeddings




class DocumentProcessor:
    def __init__(self, chunk_size=1000, chunk_overlap=200):
        embeddings = OpenAIEmbeddings()
                
        self.text_splitter = SemanticChunker(
                    embeddings,
                    breakpoint_threshold_type="percentile",  
                    breakpoint_threshold_amount=95, 
        )
        print("[INFO] Using semantic chunking")

    def load_any_file(self, path):
        ext = os.path.splitext(path)[1].lower()
        print(f"[INFO] Processing file: {path} (type: {ext})")

        if not os.path.exists(path):
            print(f"[ERROR] File does not exist: {path}")
            return []
        
        if not os.access(path, os.R_OK):
            print(f"[ERROR] File is not readable: {path}")
            return []

        try:
            if ext == ".pdf":
                return self._load_pdf(path)
            elif ext == ".docx":
                return self._load_docx(path)
            elif ext == ".pptx":
                return self._load_pptx(path)
            elif ext in [".xls", ".xlsx"]:
                return self._load_excel(path)
            elif ext == ".txt":
                return self._load_txt(path)
            else:
                print(f"[ERROR] Unsupported file type: {ext}")
                return []
        except Exception as e:
            print(f"[ERROR] Failed to process {path}: {str(e)}")
            return []

    def _load_pdf(self, path):
        """Load PDF using PyPDFLoader (digital PDFs only)"""
        try:
            print("[DEBUG] Loading PDF...")
            docs = PyPDFLoader(path).load()
            # Add metadata to each doc
            for doc in docs:
                doc.metadata.setdefault("source", path)
                doc.metadata.setdefault("doc_type", os.path.splitext(os.path.basename(path))[0])
            if docs and any(doc.page_content.strip() for doc in docs):
                print(f"[INFO] Successfully loaded PDF: {len(docs)} pages")
                return docs
            else:
                print("[WARN] PDF appears to be scanned or has no extractable text")
                return []
        except Exception as e:
            print(f"[ERROR] Failed to load PDF: {e}")
            return []

    def _load_docx(self, path):
        """Load DOCX using python-docx"""
        try:
            from docx import Document as DocxDocument
            print("[DEBUG] Loading DOCX...")
            
            doc = DocxDocument(path)
            
            # Extract text from paragraphs
            text_content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            full_text = "\n".join(text_content)
            
            if full_text.strip():
                docs = [Document(
                    page_content=full_text,
                    metadata={"source": path, "doc_type": os.path.splitext(os.path.basename(path))[0]}
                    
                )]
                print(f"[INFO] Successfully loaded DOCX: {len(full_text)} characters")
                return docs
            else:
                print("[WARN] No text content found in DOCX file")
                return []
                
        except ImportError:
            print("[ERROR] python-docx not installed. Install with: pip install python-docx")
            return []
        except Exception as e:
            print(f"[ERROR] Failed to load DOCX: {e}")
            return []

    def _load_pptx(self, path):
        """Load PPTX using python-pptx"""
        try:
            from pptx import Presentation
            print("[DEBUG] Loading PPTX...")
            
            prs = Presentation(path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"--- Slide {slide_num} ---")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract text from tables in slides
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_content.extend(slide_text)
                    text_content.append("")  # Empty line between slides
            
            full_text = "\n".join(text_content)
            
            if full_text.strip():
                docs = [Document(
                    page_content=full_text,
                    metadata={"source": path, "doc_type": os.path.splitext(os.path.basename(path))[0], "slides": len(prs.slides)}
                )]
                print(f"[INFO] Successfully loaded PPTX: {len(prs.slides)} slides, {len(full_text)} characters")
                return docs
            else:
                print("[WARN] No text content found in PPTX file")
                return []
                
        except ImportError:
            print("[ERROR] python-pptx not installed. Install with: pip install python-pptx")
            return []
        except Exception as e:
            print(f"[ERROR] Failed to load PPTX: {e}")
            return []

    def _load_excel(self, path):
        """Load Excel using pandas and openpyxl"""
        try:
            import pandas as pd
            print("[DEBUG] Loading Excel...")
            
            # Read all sheets
            excel_file = pd.ExcelFile(path)
            text_content = []
            
            for sheet_name in excel_file.sheet_names:
                print(f"[DEBUG] Processing sheet: {sheet_name}")
                df = pd.read_excel(path, sheet_name=sheet_name)
                
                # Add sheet header
                text_content.append(f"=== Sheet: {sheet_name} ===")
                
                # Convert DataFrame to text
                if not df.empty:
                    # Add column headers
                    text_content.append(" | ".join(str(col) for col in df.columns))
                    text_content.append("-" * 50)
                    
                    # Add rows
                    for _, row in df.iterrows():
                        row_text = " | ".join(str(val) for val in row.values if pd.notna(val))
                        if row_text.strip():
                            text_content.append(row_text)
                else:
                    text_content.append("(Empty sheet)")
                
                text_content.append("")  # Empty line between sheets
            
            full_text = "\n".join(text_content)
            
            if full_text.strip():
                docs = [Document(
                    page_content=full_text,
                    metadata={"source": path, "doc_type": os.path.splitext(os.path.basename(path))[0], "sheets": len(excel_file.sheet_names)}
                )]
                print(f"[INFO] Successfully loaded Excel: {len(excel_file.sheet_names)} sheets, {len(full_text)} characters")
                return docs
            else:
                print("[WARN] No text content found in Excel file")
                return []
                
        except ImportError:
            print("[ERROR] pandas or openpyxl not installed. Install with: pip install pandas openpyxl")
            return []
        except Exception as e:
            print(f"[ERROR] Failed to load Excel: {e}")
            return []

    def _load_txt(self, path):
        """Load plain text files"""
        try:
            print("[DEBUG] Loading TXT file...")
            
            # Try different encodings
            encodings = ['utf-8', 'utf-8-sig', 'latin1', 'cp1252', 'iso-8859-1']
            
            for encoding in encodings:
                try:
                    with open(path, 'r', encoding=encoding) as file:
                        content = file.read()
                    print(f"[DEBUG] Successfully read with {encoding} encoding")
                    break
                except UnicodeDecodeError:
                    continue
            else:
                print("[ERROR] Could not decode text file with any supported encoding")
                return []
            
            if content.strip():
                docs = [Document(
                    page_content=content,
                    metadata={"source": path, "doc_type": os.path.splitext(os.path.basename(path))[0]}
                )]
                print(f"[INFO] Successfully loaded TXT: {len(content)} characters")
                return docs
            else:
                print("[WARN] Text file is empty")
                return []
                
        except Exception as e:
            print(f"[ERROR] Failed to load TXT: {e}")
            return []

    def load_directory(self, folder):
        """Load all supported files from a directory"""
        all_docs = []
        print(f"[INFO] Scanning directory: {folder}")
        
        if not os.path.exists(folder):
            print(f"[ERROR] Directory does not exist: {folder}")
            return []
        
        supported_extensions = {'.pdf', '.docx', '.pptx', '.xls', '.xlsx', '.txt'}
        
        for root, _, files in os.walk(folder):
            for file in files:
                file_ext = os.path.splitext(file)[1].lower()
                if file_ext in supported_extensions:
                    file_path = os.path.join(root, file)
                    print(f"[DEBUG] Processing: {file_path}")
                    docs = self.load_any_file(file_path)
                    for doc in docs:
                        doc.metadata.setdefault("source_folder", root)
                        doc.metadata.setdefault("filename", file)
                    all_docs.extend(docs)
                else:
                    print(f"[SKIP] Unsupported file type: {file}")
        
        print(f"[INFO] Total documents loaded: {len(all_docs)}")
        return all_docs

    def split_documents(self, documents):
        """Split documents into chunks"""
        if not documents:
            print("[WARN] No documents to split")
            return []
        
        print(f"[INFO] Splitting {len(documents)} documents...")
        chunks = self.text_splitter.split_documents(documents)
        print(f"[INFO] Created {len(chunks)} chunks")
        return chunks

    def get_supported_formats(self):
        """Return list of supported file formats"""
        return ['.pdf', '.docx', '.pptx', '.xls', '.xlsx', '.txt']



